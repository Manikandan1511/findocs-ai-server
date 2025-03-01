import io
import filetype
from google.cloud import vision
from google.oauth2 import service_account
import pdfplumber
from docx import Document
from pptx import Presentation

# Load Google Vision API credentials
credentials = service_account.Credentials.from_service_account_file("service_account.json")
client = vision.ImageAnnotatorClient(credentials=credentials)

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Extracts text from images, PDFs, DOCX, and PPTX files using Google Vision API."""
    try:
        kind = filetype.guess(io.BytesIO(file_bytes))
        file_ext = kind.extension if kind else filename.split(".")[-1].lower()

        # **1. Process Handwritten & Printed Text in Images**
        if file_ext in ["jpg", "jpeg", "png", "tiff"]:
            print("✅ Processing image with Google Vision OCR")

            image = vision.Image(content=file_bytes)
            response = client.text_detection(image=image)
            extracted_text = response.text_annotations[0].description if response.text_annotations else "⚠️ No text detected in image."
            return extracted_text.strip()

        # **2. Process PDFs**
        elif file_ext == "pdf":
            print("✅ Extracting text from PDF")
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                extracted_text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
            return extracted_text if extracted_text else "No text found in PDF."

        # **3. Process DOCX (Microsoft Word)**
        elif file_ext == "docx":
            print("✅ Extracting text from DOCX")
            doc = Document(io.BytesIO(file_bytes))
            extracted_text = "\n".join([p.text for p in doc.paragraphs])
            return extracted_text if extracted_text else "No text found in DOCX."

        # **4. Process PPTX (Microsoft PowerPoint)**
        elif file_ext in ["ppt", "pptx"]:
            print("✅ Extracting text from PPTX")
            ppt = Presentation(io.BytesIO(file_bytes))
            extracted_text = "\n".join([slide.shapes.title.text for slide in ppt.slides if slide.shapes.title])
            return extracted_text if extracted_text else "No text found in PPTX."

        # **5. Unsupported Format**
        else:
            return f"❌ Unsupported file format: {file_ext}"

    except Exception as e:
        print(f"❌ OCR Error: {e}")
        return str(e)